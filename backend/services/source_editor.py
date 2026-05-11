from __future__ import annotations

import json
import re
from html import escape
from datetime import datetime
from pathlib import Path
from typing import Iterable

import pdfplumber
from docx import Document
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer
from sqlalchemy import select
from sqlalchemy.orm import Session

from .. import models
from ..config import settings
from ..llm import LLMResult
from ..rag import retriever
from .adaptation import _resolve_provider

ALLOWED_EXTENSIONS = {".docx", ".pptx", ".pdf"}
EDIT_DIR = settings.uploads_dir / "lesson_edits"
EDIT_DIR.mkdir(parents=True, exist_ok=True)


def _norm(text: str) -> set[str]:
    return {t for t in re.findall(r"[a-z0-9]+", text.lower()) if len(t) > 2}


def _safe_filename(name: str) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9._-]+", "_", name).strip("._")
    return cleaned or "lesson_edit"


def _resolve_source_path(source_path: str) -> Path:
    base = settings.sample_lessons_dir.resolve()
    candidate = (base / source_path).resolve()
    if base not in candidate.parents and candidate != base:
        raise LookupError("source file must be inside Sample Lessons")
    if not candidate.exists() or candidate.suffix.lower() not in ALLOWED_EXTENSIONS:
        raise LookupError("source file not found or unsupported")
    return candidate


def source_files_for_lesson(lesson: models.Lesson) -> list[dict]:
    lesson_tokens = _norm(lesson.title)
    files: list[dict] = []
    for path in settings.sample_lessons_dir.rglob("*"):
        if not path.is_file() or path.suffix.lower() not in ALLOWED_EXTENSIONS:
            continue
        rel = path.relative_to(settings.sample_lessons_dir)
        rel_tokens = _norm(str(rel))
        score = len(lesson_tokens & rel_tokens)
        if score == 0:
            continue
        files.append(
            {
                "source_path": rel.as_posix(),
                "filename": path.name,
                "file_type": path.suffix.lower().lstrip("."),
                "size_bytes": path.stat().st_size,
                "_score": score,
            }
        )
    files.sort(key=lambda f: (-f["_score"], f["filename"].lower()))
    for item in files:
        item.pop("_score", None)
    return files


def _docx_blocks(path: Path) -> tuple[Document, list[tuple[str, object]]]:
    doc = Document(str(path))
    blocks: list[tuple[str, object]] = []
    for p in doc.paragraphs:
        if p.text.strip():
            blocks.append((p.text, p))
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    if p.text.strip():
                        blocks.append((p.text, p))
    return doc, blocks


def _pptx_blocks(path: Path) -> tuple[Presentation, list[tuple[str, object]]]:
    prs = Presentation(str(path))
    blocks: list[tuple[str, object]] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                blocks.append((shape.text, shape))
    return prs, blocks


def _pdf_blocks(path: Path) -> list[tuple[str, object]]:
    blocks: list[tuple[str, object]] = []
    with pdfplumber.open(str(path)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                blocks.append((text.strip(), {"page": i}))
    return blocks


def _split_text_for_lengths(text: str, lengths: list[int]) -> list[str]:
    if not lengths:
        return []
    parts: list[str] = []
    pos = 0
    for i, length in enumerate(lengths):
        if i == len(lengths) - 1:
            parts.append(text[pos:])
        else:
            parts.append(text[pos : pos + length])
            pos += length
    return parts


def _set_docx_paragraph_text_preserving_runs(paragraph, new_text: str) -> None:
    # Editing the existing w:t nodes preserves paragraph style, run formatting,
    # and hyperlink XML containers far better than assigning paragraph.text.
    text_nodes = [
        node
        for node in paragraph._p.iter()
        if node.tag == "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t"
    ]
    if not text_nodes:
        paragraph.text = new_text
        return
    lengths = [len(node.text or "") for node in text_nodes]
    if not any(lengths):
        lengths = [0 for _ in text_nodes]
        lengths[0] = len(new_text)
    for node, part in zip(text_nodes, _split_text_for_lengths(new_text, lengths)):
        node.text = part


def _set_pptx_shape_text_preserving_runs(shape, new_text: str) -> None:
    runs = []
    for paragraph in shape.text_frame.paragraphs:
        runs.extend(paragraph.runs)
    if not runs:
        shape.text = new_text
        return
    lengths = [len(run.text or "") for run in runs]
    if not any(lengths):
        lengths = [0 for _ in runs]
        lengths[0] = len(new_text)
    for run, part in zip(runs, _split_text_for_lengths(new_text, lengths)):
        run.text = part


def _chunks(items: list[tuple[int, str]], max_chars: int = 7000) -> Iterable[list[tuple[int, str]]]:
    chunk: list[tuple[int, str]] = []
    size = 0
    for item_id, text in items:
        item_size = len(text)
        if chunk and size + item_size > max_chars:
            yield chunk
            chunk = []
            size = 0
        chunk.append((item_id, text))
        size += item_size
    if chunk:
        yield chunk


_JSON_FENCE = re.compile(r"```(?:json)?\s*(\{.*?\})\s*```", re.DOTALL)


def _parse_items(text: str) -> dict[int, str]:
    candidate = text.strip()
    match = _JSON_FENCE.search(candidate)
    if match:
        candidate = match.group(1)
    else:
        first = candidate.find("{")
        last = candidate.rfind("}")
        if first != -1 and last != -1:
            candidate = candidate[first : last + 1]
    data = json.loads(candidate)
    return {int(item["id"]): str(item["text"]) for item in data.get("items", [])}


def _kb_context(db: Session, lesson: models.Lesson, cluster_id: int | None, kb_ids: list[int]) -> str:
    if not kb_ids:
        return ""
    kb_rows = list(
        db.execute(select(models.KnowledgeBase).where(models.KnowledgeBase.kb_id.in_(kb_ids))).scalars()
    )
    specs = [{"kb_id": kb.kb_id, "kb_name": kb.kb_name, "category": kb.category} for kb in kb_rows]
    cluster = db.get(models.StudentCluster, cluster_id) if cluster_id else None
    query = " ".join(
        filter(
            None,
            [
                lesson.title,
                lesson.cs_topic,
                lesson.objectives,
                cluster.cluster_name if cluster else None,
                cluster.cluster_description if cluster else None,
            ],
        )
    )
    chunks = retriever.retrieve_for_lesson(query, specs, top_k_per_kb=2)
    lines: list[str] = []
    if cluster:
        lines.append(f"Learner profile: {cluster.cluster_name} - {cluster.cluster_description or ''}")
    for chunk in chunks:
        lines.append(f"[KB #{chunk.kb_id} {chunk.kb_name} - {chunk.section_title}]\n{chunk.text}")
    return "\n\n".join(lines)


def _rewrite_texts(
    db: Session,
    *,
    teacher: models.Teacher,
    lesson: models.Lesson,
    source_filename: str,
    instruction: str,
    texts: list[str],
    cluster_id: int | None,
    kb_ids: list[int],
) -> tuple[list[str], LLMResult]:
    provider = _resolve_provider(db, teacher.teacher_id)
    indexed = [(i, text) for i, text in enumerate(texts)]
    rewritten = list(texts)
    last_result: LLMResult | None = None
    kb_context = _kb_context(db, lesson, cluster_id, kb_ids)
    system = (
        "You edit teacher lesson source files. Preserve meaning, classroom usefulness, numbering, "
        "placeholders, URLs, hyperlink labels, and references unless the teacher explicitly asks to change them. "
        "Follow the teacher instruction. Return JSON only."
    )
    for batch in _chunks(indexed):
        user = {
            "lesson": {
                "title": lesson.title,
                "grade_level": lesson.grade_level,
                "cs_topic": lesson.cs_topic,
                "objectives": lesson.objectives,
            },
            "source_file": source_filename,
            "teacher_instruction": instruction,
            "rag_context": kb_context,
            "items": [{"id": item_id, "text": text} for item_id, text in batch],
            "required_output": {
                "items": [
                    {
                        "id": "same numeric id from input",
                        "text": "edited replacement text only; do not add explanations",
                    }
                ]
            },
        }
        result = provider.generate(system=system, user=json.dumps(user, ensure_ascii=False), max_tokens=8192)
        last_result = result
        try:
            parsed = _parse_items(result.text)
        except Exception as e:
            raise RuntimeError("LLM returned an unreadable file-edit response") from e
        for item_id, new_text in parsed.items():
            if 0 <= item_id < len(rewritten) and new_text.strip():
                rewritten[item_id] = new_text.strip()
    if last_result is None:
        raise RuntimeError("source file did not contain editable text")
    return rewritten, last_result


def _write_pdf(out_path: Path, title: str, blocks: list[str]) -> None:
    styles = getSampleStyleSheet()
    doc = SimpleDocTemplate(str(out_path), pagesize=letter, title=title)
    story = [Paragraph(title, styles["Title"]), Spacer(1, 12)]
    for i, text in enumerate(blocks):
        if i:
            story.append(PageBreak())
        for para in text.splitlines():
            if para.strip():
                story.append(Paragraph(escape(para.strip()), styles["BodyText"]))
                story.append(Spacer(1, 6))
    doc.build(story)


def edit_source_file(
    db: Session,
    *,
    teacher: models.Teacher,
    lesson_id: int,
    source_path: str,
    instruction: str,
    cluster_id: int | None,
    kb_ids: list[int],
) -> dict:
    lesson = db.get(models.Lesson, lesson_id)
    if not lesson:
        raise LookupError("lesson not found")
    source = _resolve_source_path(source_path)
    rel = source.relative_to(settings.sample_lessons_dir).as_posix()
    if rel not in {f["source_path"] for f in source_files_for_lesson(lesson)}:
        raise LookupError("source file is not linked to this lesson")

    suffix = source.suffix.lower()
    stamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    out_name = _safe_filename(f"{source.stem}_ai_edit_{stamp}{suffix}")
    out_path = EDIT_DIR / out_name

    if suffix == ".docx":
        doc, blocks = _docx_blocks(source)
        texts, targets = [b[0] for b in blocks], [b[1] for b in blocks]
        rewritten, _result = _rewrite_texts(
            db,
            teacher=teacher,
            lesson=lesson,
            source_filename=source.name,
            instruction=instruction,
            texts=texts,
            cluster_id=cluster_id,
            kb_ids=kb_ids,
        )
        for target, new_text in zip(targets, rewritten):
            _set_docx_paragraph_text_preserving_runs(target, new_text)
        doc.save(str(out_path))
    elif suffix == ".pptx":
        prs, blocks = _pptx_blocks(source)
        texts, targets = [b[0] for b in blocks], [b[1] for b in blocks]
        rewritten, _result = _rewrite_texts(
            db,
            teacher=teacher,
            lesson=lesson,
            source_filename=source.name,
            instruction=instruction,
            texts=texts,
            cluster_id=cluster_id,
            kb_ids=kb_ids,
        )
        for target, new_text in zip(targets, rewritten):
            _set_pptx_shape_text_preserving_runs(target, new_text)
        prs.save(str(out_path))
    elif suffix == ".pdf":
        blocks = _pdf_blocks(source)
        texts = [b[0] for b in blocks]
        rewritten, _result = _rewrite_texts(
            db,
            teacher=teacher,
            lesson=lesson,
            source_filename=source.name,
            instruction=instruction,
            texts=texts,
            cluster_id=cluster_id,
            kb_ids=kb_ids,
        )
        _write_pdf(out_path, f"{lesson.title} - AI Edited Copy", rewritten)
    else:
        raise LookupError("unsupported source file type")

    return {
        "filename": out_name,
        "file_type": suffix.lstrip("."),
        "download_url": f"/api/lesson-file-edits/{out_name}",
        "note": (
            "Original source file was not changed. DOCX/PPTX copies preserve existing text runs, hyperlinks, and basic document structure where possible; "
            "PDF copies are regenerated as text-based PDFs rather than exact-layout edits."
        ),
    }


def edited_file_path(filename: str) -> Path:
    safe = Path(filename).name
    path = (EDIT_DIR / safe).resolve()
    if EDIT_DIR.resolve() not in path.parents or not path.exists():
        raise LookupError("edited file not found")
    return path
